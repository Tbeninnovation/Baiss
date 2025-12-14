# here is the output format of the json at the end of the file
# [
#     {
#         "page": 1,
#         "data": "text"
#     },
#     {
#         "page": 2,
#         "data": "text"
#     }
# ]

import logging
from io import BytesIO
from pptx import Presentation
from baiss_sdk.parsers import extract_chunks

class PPTExtractor:
    def __init__(self):
        super().__init__()


    def get_bulleted_text(self, paragraph):
        """
        Get bulleted text from a paragraph.
        """
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  '* paragraph.level}.{paragraph.text}"
        else:
            return paragraph.text

    def extract(self, shape):
        """
        Extract text from a shape.
        """
        try:
            # Initialize output at the beginning
            output = {}

            # First try to get text content
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.get_bulleted_text(paragraph))
                output["text"] = "\n".join(texts)
                output["type"] = "text"
                return output

            # Safely get shape_type
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                # If shape_type is not available, try to get text content
                if hasattr(shape, 'text'):
                    output["text"] = shape.text.strip()
                    output["type"] = "text"
                    return output
                return ""

            # Handle table
            if shape_type == 19:
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    rows.append("; ".join([tb.cell(
                        0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
                output["text"] = "\n".join(rows)
                output["type"] = "table"
                return output

            # Handle group shape
            if shape_type == 6:
                texts = []
                for p in sorted(shape.shapes, key=lambda x: (x.top // 10, x.left)):
                    t = self.extract(p)
                    if t:
                        if isinstance(t, dict):
                            texts.append(t["text"])
                        else:
                            texts.append(t)
                output["text"] = "\n".join(texts)
                output["type"] = "group"
                return output

            # For any other shape type, try to get text if available
            if hasattr(shape, 'text') and shape.text:
                output["text"] = shape.text.strip()
                output["type"] = "text"
                return output

            return ""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return ""

    def get_total_page(self, fnp):
        """
        Get total page of a pptx file.
        """
        ppt = Presentation(fnp) if isinstance(
            fnp, str) else Presentation(
            BytesIO(fnp))
        return len(ppt.slides)

    def call_all_pages(self, fnp):
        """
        Call all pages of a pptx file.
        """
        ppt = Presentation(fnp) if isinstance(
            fnp, str) else Presentation(
            BytesIO(fnp))
        output = []

        total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            one_page_output = {}
            if i >= total_page:
                break
            texts = []
            for shape in sorted(
                    slide.shapes, key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left)):
                try:
                    txt = self.extract(shape)
                    if txt:
                        if isinstance(txt, dict):
                            texts.append(txt["text"])
                        else:
                            texts.append(txt)
                except Exception as e:
                    logging.exception(e)
            one_page_output["data"]   = "\n".join(texts)
            one_page_output["page"]   = i + 1  # Start from slide 1 instead of 0
            one_page_output["chunks"] = extract_chunks( one_page_output["data"] )
            output.append(one_page_output)

        return output

if __name__ == "__main__":
    pass
